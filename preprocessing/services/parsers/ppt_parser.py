import io
from typing import Dict, Any
from datasets import Dataset
from .base import BaseParser
from pptx import Presentation

class PPTParser(BaseParser):
    """Parser for PowerPoint presentations"""
    
    def parse(self, file_stream: io.BytesIO) -> Dataset:
        """Parse a PPTX file from an in-memory stream.
        
        Args:
            file_stream: An in-memory binary stream (e.g., io.BytesIO)
                containing the PPTX file content.
            
        Returns:
            A Hugging Face Dataset with 'slide_number' and 'text' columns.
        """
        try:
            prs = Presentation(file_stream)

    
            data = []

            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:                 
                        for paragraph in shape.text_frame.paragraphs :
                            if paragraph.text:
                                slide_text.append(paragraph.text.strip())

                if slide_text:
                    full_slide_text = "\\n".join(slide_text)
                    data.append({"slide_number": i+1, "text": full_slide_text})

            return Dataset.from_list(data)
        except Exception as e:
            print(f"Error parsing PPTX file {file_path}: {e}")
            return Dataset.from_list([])
        

